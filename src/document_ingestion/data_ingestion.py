import os
import sys
import json
import uuid
import hashlib
import shutil
import pandas as pd
from pathlib import Path
from typing import Iterable, List, Optional, Dict, Any, Union
import fitz  # PyMuPDF
from langchain.schema import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from utils.model_loader import ModelLoader
from logger import GLOBAL_LOGGER as log
from exception.custom_exception import DocumentPortalException
from utils.file_io import generate_session_id, save_uploaded_files
from utils.document_ops import load_documents, concat_for_analysis, concat_for_comparison

# Import required libraries for document processing
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    log.warning("python-docx not available. Install with: pip install python-docx")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    log.warning("python-pptx not available. Install with: pip install python-pptx")

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    log.warning("openpyxl not available. Install with: pip install openpyxl")

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".ppt", ".pptx", 
    ".xlsx", ".xls", ".csv", ".json", ".rtf"}

class DocHandler:
    """
    Universal document save + read for analysis - supports all document types.
    """
    def __init__(self, data_dir: Optional[str] = None, session_id: Optional[str] = None):
        self.data_dir = data_dir or os.getenv("DATA_STORAGE_PATH", os.path.join(os.getcwd(), "data", "document_analysis"))
        self.session_id = session_id or generate_session_id("session")
        self.session_path = os.path.join(self.data_dir, self.session_id)
        os.makedirs(self.session_path, exist_ok=True)
        log.info("DocHandler initialized", session_id=self.session_id, session_path=self.session_path)

    def save_pdf(self, uploaded_file) -> str:
        """Legacy method for backward compatibility"""
        return self.save_document(uploaded_file)

    def save_document(self, uploaded_file) -> str:
        """Save any supported document type"""
        try:
            filename = os.path.basename(uploaded_file.name)
            file_ext = Path(filename).suffix.lower()
            
            if file_ext not in SUPPORTED_EXTENSIONS:
                raise ValueError(f'Unsupported file type: {file_ext}. Supported: {", ".join(SUPPORTED_EXTENSIONS)}')
            
            save_path = os.path.join(self.session_path, filename)
            with open(save_path, "wb") as f:
                if hasattr(uploaded_file, "read"):
                    f.write(uploaded_file.read())
                else:
                    f.write(uploaded_file.getbuffer())
                    
            log.info("Document saved successfully", file=filename, type=file_ext, save_path=save_path, session_id=self.session_id)
            return save_path
        except Exception as e:
            log.error("Failed to save document", error=str(e), session_id=self.session_id)
            raise DocumentPortalException(f"Failed to save document: {str(e)}", e) from e

    def read_document(self, doc_path: str) -> str:
        """Read any supported document type"""
        try:
            file_path = Path(doc_path)
            file_ext = file_path.suffix.lower()
            
            if file_ext == '.pdf':
                return self._read_pdf(doc_path)
            elif file_ext == '.docx':
                return self._read_docx(doc_path)
            elif file_ext in ['.ppt', '.pptx']:
                return self._read_ppt(doc_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._read_excel(doc_path)
            elif file_ext == '.csv':
                return self._read_csv(doc_path)
            elif file_ext in ['.txt', '.md']:
                return self._read_text(doc_path)
            elif file_ext == '.json':
                return self._read_json(doc_path)
            elif file_ext == '.rtf':
                return self._read_rtf(doc_path)
            else:
                raise ValueError(f"Unsupported file extension: {file_ext}")
                
        except Exception as e:
            log.error('Failed to read document', error=str(e), doc_path=doc_path, session_id=self.session_id)
            raise DocumentPortalException(f'Could not process document: {doc_path}', e) from e

    def _read_pdf(self, pdf_path: str) -> str:
        """Read PDF documents using PyMuPDF"""
        try:
            text_chunks = []
            with fitz.open(pdf_path) as doc:
                for page_num in range(doc.page_count):
                    page = doc.load_page(page_num)
                    page_text = page.get_text()
                    if page_text.strip():  # Only add non-empty pages
                        text_chunks.append(f"\n--- Page {page_num + 1} ---\n{page_text}")
            
            text = "\n".join(text_chunks)
            log.info("PDF read successfully", pdf_path=pdf_path, pages=len(text_chunks))
            return text
        except Exception as e:
            raise DocumentPortalException(f"Error reading PDF: {pdf_path}", e) from e

    def _read_docx(self, docx_path: str) -> str:
        """Read Word documents using python-docx"""
        if not DOCX_AVAILABLE:
            raise DocumentPortalException("python-docx library not available. Install with: pip install python-docx", None)
        
        try:
            doc = DocxDocument(docx_path)
            text_parts = []
            
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                
                if table_text:
                    text_parts.append("\n--- Table ---\n" + "\n".join(table_text))
            
            content = "\n\n".join(text_parts)
            log.info("DOCX read successfully", path=docx_path, paragraphs=len(doc.paragraphs), tables=len(doc.tables))
            return content
            
        except Exception as e:
            raise DocumentPortalException(f"Error reading DOCX: {docx_path}", e) from e

    def _read_ppt(self, ppt_path: str) -> str:
        """Read PowerPoint documents using python-pptx"""
        if not PPTX_AVAILABLE:
            raise DocumentPortalException("python-pptx library not available. Install with: pip install python-pptx", None)
        
        try:
            prs = Presentation(ppt_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables in slides
                    if shape.has_table:
                        table_text = ["Table in slide:"]
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_data.append(cell.text.strip())
                            if row_data:
                                table_text.append(" | ".join(row_data))
                        slide_text.extend(table_text)
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_parts.append("\n".join(slide_text))
            
            content = "\n\n".join(text_parts)
            log.info("PPT read successfully", path=ppt_path, slides=len(prs.slides))
            return content
            
        except Exception as e:
            raise DocumentPortalException(f"Error reading PPT: {ppt_path}", e) from e

    def _read_excel(self, excel_path: str) -> str:
        """Read Excel documents using pandas and openpyxl"""
        if not EXCEL_AVAILABLE:
            raise DocumentPortalException("openpyxl library not available. Install with: pip install openpyxl", None)
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(excel_path)
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_path, sheet_name=sheet_name)
                
                if not df.empty:
                    sheet_text = [f"--- Sheet: {sheet_name} ---"]
                    
                    # Add column headers
                    headers = " | ".join(str(col) for col in df.columns)
                    sheet_text.append(f"Headers: {headers}")
                    
                    # Add data (limit to first 100 rows to avoid huge output)
                    max_rows = min(100, len(df))
                    for idx, row in df.head(max_rows).iterrows():
                        row_data = " | ".join(str(val) for val in row.values if pd.notna(val))
                        if row_data.strip():
                            sheet_text.append(row_data)
                    
                    if len(df) > max_rows:
                        sheet_text.append(f"... and {len(df) - max_rows} more rows")
                    
                    text_parts.append("\n".join(sheet_text))
            
            content = "\n\n".join(text_parts)
            log.info("Excel read successfully", path=excel_path, sheets=len(excel_file.sheet_names))
            return content
            
        except Exception as e:
            raise DocumentPortalException(f"Error reading Excel: {excel_path}", e) from e

    def _read_csv(self, csv_path: str) -> str:
        """Read CSV documents using pandas"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(csv_path, encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                raise ValueError("Could not decode CSV with any common encoding")
            
            if df.empty:
                return "CSV file is empty"
            
            text_parts = ["--- CSV Data ---"]
            
            # Add column headers
            headers = " | ".join(str(col) for col in df.columns)
            text_parts.append(f"Headers: {headers}")
            
            # Add data (limit to first 100 rows)
            max_rows = min(100, len(df))
            for idx, row in df.head(max_rows).iterrows():
                row_data = " | ".join(str(val) for val in row.values if pd.notna(val))
                if row_data.strip():
                    text_parts.append(row_data)
            
            if len(df) > max_rows:
                text_parts.append(f"... and {len(df) - max_rows} more rows")
            
            content = "\n".join(text_parts)
            log.info("CSV read successfully", path=csv_path, rows=len(df), columns=len(df.columns))
            return content
            
        except Exception as e:
            raise DocumentPortalException(f"Error reading CSV: {csv_path}", e) from e

    def _read_text(self, text_path: str) -> str:
        """Read plain text and markdown files"""
        try:
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            content = None
            
            for encoding in encodings:
                try:
                    with open(text_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                raise ValueError("Could not decode text file with any common encoding")
                
            log.info("Text file read successfully", path=text_path, length=len(content))
            return content
            
        except Exception as e:
            raise DocumentPortalException(f"Error reading text file: {text_path}", e) from e

    def _read_json(self, json_path: str) -> str:
        """Read JSON files"""
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Format JSON nicely for analysis
            content = json.dumps(data, indent=2, ensure_ascii=False)
            log.info("JSON file read successfully", path=json_path)
            return content
            
        except Exception as e:
            raise DocumentPortalException(f"Error reading JSON: {json_path}", e) from e

    def _read_rtf(self, rtf_path: str) -> str:
        """Read RTF files - basic implementation"""
        try:
            # For now, read RTF as plain text (not ideal but functional)
            # For better RTF support, consider using striprtf library
            with open(rtf_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # Basic RTF cleaning - remove common RTF tags
            import re
            content = re.sub(r'\\[a-z]+\d*\s?', '', content)  # Remove RTF control words
            content = re.sub(r'[{}]', '', content)  # Remove braces
            content = '\n'.join(line.strip() for line in content.split('\n') if line.strip())
            
            log.info("RTF file read successfully", path=rtf_path)
            return content
            
        except Exception as e:
            raise DocumentPortalException(f"Error reading RTF: {rtf_path}", e) from e